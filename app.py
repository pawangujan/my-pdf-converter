import streamlit as st
from pdf2image import convert_from_bytes
from pptx import Presentation
from pptx.util import Inches
import io
import os

st.title("🛡️ Secure PDF to PPT Converter")

# 1. UPDATE THIS TO YOUR EXACT 'BIN' FOLDER
POPPLER_PATH = r"C:\Users\ACER\Downloads\Release-25.12.0-0\poppler-25.12.0\Library\bin"

uploaded_file = st.file_uploader("Upload PDF", type="pdf")

if uploaded_file and st.button("Convert Now"):
    # Check if the path actually exists first
    if not os.path.exists(POPPLER_PATH):
        st.error(f"❌ FOLDER NOT FOUND: {POPPLER_PATH}")
        st.info("Please check the 'POPPLER_PATH' in your code. It must match your folder exactly.")
    else:
        try:
            with st.spinner("Converting..."):
                images = convert_from_bytes(uploaded_file.read(), dpi=100, poppler_path=POPPLER_PATH)
                
                prs = Presentation()
                for img in images:
                    # Get size and scale to fit PPT limits
                    w_px, h_px = img.size
                    slide_w, slide_h = w_px/100, h_px/100
                    
                    if slide_w > 50 or slide_h > 50:
                        scale = 50 / max(slide_w, slide_h)
                        slide_w, slide_h = slide_w * scale, slide_h * scale

                    prs.slide_width = Inches(slide_w)
                    prs.slide_height = Inches(slide_h)
                    
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    img_io = io.BytesIO()
                    img.save(img_io, format='JPEG', quality=80)
                    slide.shapes.add_picture(img_io, 0, 0, width=Inches(slide_w), height=Inches(slide_h))
                
                ppt_io = io.BytesIO()
                prs.save(ppt_io)
                st.success("✅ Success!")
                st.download_button("📥 Download PPTX", ppt_io.getvalue(), "converted.pptx")
        except Exception as e:
            st.error(f"Error: {e}")